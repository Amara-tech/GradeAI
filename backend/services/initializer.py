from typing import List
import pypdf
import pytesseract
from PIL import Image
import os
import io
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader
from .router import DocumentIngestionRouter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_experimental.text_splitter import SemanticChunker

class GraphRAGInitializer:
    """
    The orchestrator for the backend's ingestion phase.
    Accepts raw materials, routes them based on structural characteristics,
    and returns a uniform list of LangChain Document objects.
    """
    
    def __init__(self):
        # Initialize our router with the internal class processing methods
        self.router = DocumentIngestionRouter(
            text_pdf_handler=self._load_digital_pdf,
            scanned_pdf_handler=self._load_scanned_pdf,
            slide_image_handler=self._load_slides_and_images,
            docx_handler=self._load_docx
        )
        # This will hold our unified documents during processing
        self._processed_documents: List[Document] = []

    def process_material(self, file_path: str) -> List[Document]:
        """
        Public entry point. Pass any course material file path here.
        Returns a clean list of LangChain Documents.
        """
        self._processed_documents.clear()  # Reset buffer for new file
        
        print(f"[Initializer] Processing: {file_path}")
        self.router.route_document(file_path)
        
        return self._processed_documents

    def _load_digital_pdf(self, file_path: str) -> None:
        """
        Handles native digital text extraction using standard LangChain loaders.
        """
        
        loader = PyPDFLoader(file_path)
        # Extends our buffer with standard LangChain Document objects
        self._processed_documents.extend(loader.load())

    def _load_scanned_pdf(self, file_path: str) -> None:
        """
            Handles image-only scans using local pypdf to extract image blobs 
            and pytesseract to convert those blobs into clean text documents.
        """
        print(f"[OCR Engine] Running Tesseract OCR on scanned PDF: {file_path}")
    
        try:
            reader = pypdf.PdfReader(file_path)
        
            for page_idx, page in enumerate(reader.pages):
                page_text_parts = []
            
                # Extract image objects embedded in the PDF page structure
                if "/Resources" in page and "/XObject" in page["/Resources"]:
                    xObject = page["/Resources"]["/XObject"].get_object()
                
                    for obj_name in xObject:
                        obj = xObject[obj_name].get_object()
                        if obj["/Subtype"] == "/Image":
                            # Convert raw PDF image bytes into a PIL Image object
                            image_bytes = obj.get_data()
                            image = Image.open(io.BytesIO(image_bytes))
                        
                            # Execute local OCR
                            extracted_text = pytesseract.image_to_string(image)
                            if extracted_text.strip():
                                page_text_parts.append(extracted_text)
            
                # If text parts were found across images on this page, unify them
                combined_page_text = "\n".join(page_text_parts).strip()
            
                if combined_page_text:
                    doc = Document(
                        page_content=combined_page_text,
                        metadata={
                            "source": file_path,
                            "page": page_idx,  # Zero-indexed page tracker
                            "type": "scanned_pdf"
                        }
                    )
                    self._processed_documents.append(doc)
                
        except Exception as e:
            print(f"[Error] Failed running local OCR on {file_path}: {str(e)}")
        
        
    def _load_docx(self, file_path: str) -> None:
        """
        Handles native Microsoft Word text extraction using Docx2txtLoader.
        """
        from langchain_community.document_loaders import Docx2txtLoader
        
        try:
            loader = Docx2txtLoader(file_path)
            # This extracts the full document text cleanly
            doc_contents = loader.load()
            
            # Since docx doesn't inherently have rigid "pages" like a PDF, 
            # we default page index to 0 for the whole document block
            for doc in doc_contents:
                doc.metadata.update({
                    "source": file_path,
                    "page": 0,
                    "type": "word_document"
                })
                
            self._processed_documents.extend(doc_contents)
        except Exception as e:
            print(f"[Error] Failed to process Word document {file_path}: {str(e)}")    
        
        
    def _load_slides_and_images(self, file_path: str) -> None:
        """
        Handles presentations (.pptx) and standalone images (.png, .jpg).
        """
        _, ext = os.path.splitext(file_path.lower())
    
        # --- Sub-Route A: Standalone Images ---
        if ext in {'.png', '.jpg', '.jpeg', '.tiff'}:
            try:
                image = Image.open(file_path)
                extracted_text = pytesseract.image_to_string(image).strip()
                if extracted_text:
                    doc = Document(
                        page_content=extracted_text,
                        metadata={
                            "source": file_path,
                            "page": 0,  # Single images default to page 0
                            "type": "raw_image"
                        }
                    )
                    self._processed_documents.append(doc)
            except Exception as e:
                print(f"[Error] Failed to process image {file_path}: {str(e)}")
            
        # --- Sub-Route B: PowerPoint Presentations ---
        elif ext in {'.pptx', '.ppt'}:
            from pptx import Presentation
        
            try:
                prs = Presentation(file_path)
            
                for slide_idx, slide in enumerate(prs.slides):
                    slide_text_parts = []
                
                    # Extract text strings cleanly from text frames within shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text_parts.append(shape.text.strip())
                        
                    combined_slide_text = "\n".join(slide_text_parts).strip()
                
                    if combined_slide_text:
                        doc = Document(
                            page_content=combined_slide_text,
                            metadata={
                                "source": file_path,
                                "page": slide_idx,  # Maps slide index seamlessly to 'page'
                                "type": "presentation_slide"
                            }
                        )
                        self._processed_documents.append(doc)
            except Exception as e:
                print(f"[Error] Failed to process PowerPoint presentation {file_path}: {str(e)}")
    
    def process_directory(self, directory_path: str) -> List[Document]:
        """
        Scans an entire directory, processes every supported course material file 
        found inside, and aggregates them into a single list of unified Documents.
        """
        if not os.path.isdir(directory_path):
            raise NotADirectoryError(f"Provided path is not a valid directory: {directory_path}")
            
        print(f"\n[Initializer] --- Starting Batch Ingestion for Directory: {directory_path} ---")
        
        # Clear previous buffers to ensure a clean exam session initialization
        self._processed_documents.clear()
        
        # Track processing statistics for the examiner's logs
        processed_count = 0
        skipped_count = 0

        # os.walk handles nested subfolders automatically if the examiner organizes by topics
        for root, _, files in os.walk(directory_path):
            for file in files:
                # Skip hidden operating system files (e.g., .DS_Store on macOS)
                if file.startswith('.'):
                    continue
                    
                file_path = os.path.join(root, file)
                _, ext = os.path.splitext(file.lower())
                
                # Check if the file matches any of our router's supported extensions
                supported_extensions = {'.pdf', '.pptx', '.ppt', '.png', '.jpg', '.jpeg', '.tiff', '.docx'}
                
                if ext in supported_extensions:
                    try:
                        # We use a secondary buffer list if we want to track files individually,
                        # but since router.route_document calls our internal handlers directly,
                        # it will seamlessly append new pages/slides to self._processed_documents.
                        self.router.route_document(file_path)
                        processed_count += 1
                    except Exception as e:
                        print(f"[Warning] Failed to process file {file}: {str(e)}")
                        skipped_count += 1
                else:
                    print(f"[Info] Skipping unsupported file type: {file}")
                    skipped_count += 1

        print(f"\n[Initializer] --- Batch Ingestion Complete ---")
        print(f"Successfully processed: {processed_count} files")
        print(f"Skipped/Failed: {skipped_count} files")
        print(f"Total unified LangChain Document segments loaded: {len(self._processed_documents)}\n")
        
        return self._processed_documents    
    
    
    def chunk_documents(self, documents: List[Document]) -> List[Document]:
        """
        Takes the unified list of raw documents and executes an optimized, 
        hierarchical character split. It respects structural elements like 
        paragraphs and line-level definitions to ensure code blocks and rubrics 
        stay together.
        """
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        
        if not documents:
            print("[Warning] No documents provided to the chunking engine.")
            return []
            
        print(f"[Chunking Engine] Running optimized structural splitter...")
        
        # We tune the chunk size to closely mimic standard question/definition lengths
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1500,       # Targets ~120-150 words per node
            chunk_overlap=150,     # Keeps continuous context sliding over splits
            separators=["\n\n", "\n", " ", ""]  # Checks paragraphs first, then line items, then words
        )
        
        # Split documents while fully retaining all your metadata attributes
        structural_chunks = text_splitter.split_documents(documents)
        
        print(f"[Chunking Engine] Partitioning complete.")
        print(f"Original segments: {len(documents)} -> Structural chunks: {len(structural_chunks)}\n")
        
        return structural_chunks